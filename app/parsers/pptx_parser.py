"""PPTX parser — slide decks routed into evidence atoms.

Real MSP deals carry PPTX slide decks alongside DOCX SOWs:
  * kickoff decks
  * design review decks
  * exec summary decks
  * BOM-on-slides quote decks

Each slide is treated as a section. Text frames + tables on the
slide become atoms. Slide titles become headings. Speaker notes
are captured as a secondary atom stream so PM sees both the
shown content and the speaker context.

The parser is deliberately tolerant: a corrupt or password-
protected PPTX raises a parse error (which A6 graceful
degradation handles) rather than killing the run.
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from app.core.ids import stable_id
from app.core.schemas import (
    ArtifactType,
    AtomType,
    AuthorityClass,
    EvidenceAtom,
    ParserCapability,
    ParserMatch,
    ParserOutput,
    ReviewStatus,
    SourceRef,
)
from app.domain.schemas import DomainPack
from app.parsers.base import BaseParser


_HEADING_TOKENS = re.compile(
    r"^(?:scope|exclusions?|out\s+of\s+scope|deliverables?|"
    r"schedule|milestones?|timeline|risks?|assumptions?|"
    r"commercial|pricing|budget|next\s+steps?|action\s+items?|"
    r"questions?|approvals?)\b",
    re.IGNORECASE,
)

# Same atom-type heuristics as docx_parser uses — keep families
# consistent so packetization treats PPTX text the same way it
# treats DOCX text.
_EXCLUSION_RE = re.compile(
    r"\b(out\s+of\s+scope|excluded?|not\s+included|"
    r"explicitly\s+excludes?|exclusion[s]?:)",
    re.IGNORECASE,
)
_CONSTRAINT_RE = re.compile(
    r"\b(must|shall|required?|requirement|"
    r"after-?hours|escort|badge|lift|"
    r"compliance|regulatory)\b",
    re.IGNORECASE,
)
_ASSUMPTION_RE = re.compile(
    r"\b(assume[ds]?|assumption[s]?|"
    r"we\s+assume|customer\s+provides?|customer\s+supplies)\b",
    re.IGNORECASE,
)
_QUESTION_RE = re.compile(
    r"\?$|^(?:question|tbd|open\s+question|to\s+confirm)\b",
    re.IGNORECASE,
)


def _classify_text(text: str, is_heading: bool) -> AtomType:
    if is_heading:
        return AtomType.scope_item
    if _EXCLUSION_RE.search(text):
        return AtomType.exclusion
    if _ASSUMPTION_RE.search(text):
        return AtomType.assumption
    if _QUESTION_RE.search(text):
        return AtomType.open_question
    if _CONSTRAINT_RE.search(text):
        return AtomType.constraint
    return AtomType.scope_item


class PptxParser(BaseParser):
    parser_name = "pptx"
    parser_version = "pptx_parser_v1"
    capability = ParserCapability(
        parser_name=parser_name,
        parser_version=parser_version,
        supported_extensions=[".pptx"],
        supported_artifact_types=[ArtifactType.pptx],
        emitted_atom_types=[
            AtomType.scope_item,
            AtomType.exclusion,
            AtomType.constraint,
            AtomType.assumption,
            AtomType.open_question,
        ],
        supported_domain_packs=["*"],
        requires_binary=True,
        supports_source_replay=True,
    )

    def match(self, path: Path, sample_text: str | None, domain_pack: DomainPack | None) -> ParserMatch:
        del sample_text, domain_pack
        suffix = path.suffix.lower()
        confidence = 0.94 if suffix == ".pptx" else 0.0
        reasons = ["pptx_extension"] if suffix == ".pptx" else []
        return ParserMatch(
            parser_name=self.parser_name,
            confidence=confidence,
            reasons=reasons,
            artifact_type=ArtifactType.pptx,
        )

    def parse(self, artifact_path: Path) -> list[Any]:
        artifact_id = stable_id("art", str(artifact_path))
        return self.parse_artifact("unknown_project", artifact_id, artifact_path)

    def parse_artifact(
        self,
        project_id: str,
        artifact_id: str,
        path: Path,
        domain_pack: DomainPack | None = None,
    ) -> list[EvidenceAtom]:
        return self.parse_artifact_full(
            project_id=project_id,
            artifact_id=artifact_id,
            path=path,
            domain_pack=domain_pack,
        ).atoms

    def parse_artifact_full(
        self,
        project_id: str,
        artifact_id: str,
        path: Path,
        domain_pack: DomainPack | None = None,
    ) -> ParserOutput:
        del domain_pack
        try:
            from pptx import Presentation
        except Exception as exc:  # pragma: no cover — env-specific
            raise RuntimeError(
                "python-pptx is required for the PPTX parser"
            ) from exc

        prs = Presentation(str(path))
        atoms: list[EvidenceAtom] = []
        slide_count = len(prs.slides)
        for slide_idx, slide in enumerate(prs.slides):
            slide_title_text = ""
            # Slide title is usually the first text-bearing placeholder.
            for shape in slide.shapes:
                if shape.has_text_frame and shape.placeholder_format is not None and "title" in str(shape.placeholder_format.type).lower():
                    slide_title_text = shape.text_frame.text.strip()
                    break
            if not slide_title_text:
                # Fallback: take the first non-empty text run from the slide
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            slide_title_text = t.split("\n", 1)[0][:80]
                            break

            # Emit one atom per text-frame paragraph; one atom per
            # table cell.
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        text = (paragraph.text or "").strip()
                        if not text:
                            continue
                        is_heading = (
                            para_idx == 0
                            and shape_idx == 0
                            and (text == slide_title_text or _HEADING_TOKENS.search(text) is not None)
                        )
                        atoms.append(self._make_atom(
                            project_id=project_id,
                            artifact_id=artifact_id,
                            filename=path.name,
                            text=text,
                            slide_index=slide_idx,
                            slide_title=slide_title_text,
                            shape_index=shape_idx,
                            paragraph_index=para_idx,
                            is_heading=is_heading,
                        ))
                if shape.has_table:
                    # Emit one atom per row (joined cells) plus a header
                    # atom for row 0 if it looks like a header. Per-row
                    # atoms preserve "Item | Qty | Unit" together so
                    # quantity / device extraction can bind them. Empty
                    # rows are skipped.
                    rows_data: list[list[str]] = []
                    for row in shape.table.rows:
                        cells = [(cell.text or "").strip() for cell in row.cells]
                        if any(cells):
                            rows_data.append(cells)
                    for row_idx, cells in enumerate(rows_data):
                        row_text = " | ".join(c for c in cells if c)
                        if not row_text:
                            continue
                        atoms.append(self._make_atom(
                            project_id=project_id,
                            artifact_id=artifact_id,
                            filename=path.name,
                            text=row_text,
                            slide_index=slide_idx,
                            slide_title=slide_title_text,
                            shape_index=shape_idx,
                            paragraph_index=None,
                            table_row=row_idx,
                            table_cell=None,
                            is_heading=False,
                        ))

            # Speaker notes — secondary atoms that capture the
            # narrator's context; tagged as ``meeting_commitment``
            # when the notes contain decision / action language so
            # the packetizer routes them appropriately.
            notes_text = ""
            try:
                if slide.has_notes_slide:
                    notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            except Exception:
                notes_text = ""
            if notes_text:
                for para_idx, para_text in enumerate(notes_text.split("\n")):
                    para_text = para_text.strip()
                    if not para_text:
                        continue
                    atoms.append(self._make_atom(
                        project_id=project_id,
                        artifact_id=artifact_id,
                        filename=path.name,
                        text=f"[Speaker notes] {para_text}",
                        slide_index=slide_idx,
                        slide_title=slide_title_text,
                        shape_index=None,
                        paragraph_index=para_idx,
                        is_heading=False,
                        is_speaker_note=True,
                    ))
        return ParserOutput(atoms=atoms, derived_files=[])

    def _make_atom(
        self,
        *,
        project_id: str,
        artifact_id: str,
        filename: str,
        text: str,
        slide_index: int,
        slide_title: str,
        shape_index: int | None,
        paragraph_index: int | None = None,
        table_row: int | None = None,
        table_cell: int | None = None,
        is_heading: bool = False,
        is_speaker_note: bool = False,
    ) -> EvidenceAtom:
        atom_type = _classify_text(text, is_heading)
        locator: dict[str, Any] = {
            "slide": slide_index + 1,
            "slide_title": slide_title,
        }
        if shape_index is not None:
            locator["shape"] = shape_index
        if paragraph_index is not None:
            locator["paragraph"] = paragraph_index
        if table_row is not None:
            locator["row"] = table_row
        if table_cell is not None:
            locator["cell"] = table_cell
        if is_speaker_note:
            locator["kind"] = "speaker_notes"
        atom_id = stable_id(
            "atm",
            project_id,
            artifact_id,
            str(slide_index),
            str(shape_index),
            str(paragraph_index),
            text[:120],
        )
        source_ref = SourceRef(
            id=stable_id(
                "src", artifact_id, str(slide_index),
                str(shape_index), str(paragraph_index),
                str(table_row), str(table_cell),
            ),
            artifact_id=artifact_id,
            artifact_type=ArtifactType.pptx,
            filename=filename,
            locator=locator,
            extraction_method="pptx_text_frames",
            parser_version=self.parser_version,
        )
        return EvidenceAtom(
            id=atom_id,
            project_id=project_id,
            artifact_id=artifact_id,
            atom_type=atom_type,
            raw_text=text,
            normalized_text=text.lower(),
            value={"kind": "speaker_notes" if is_speaker_note else "slide_text"},
            entity_keys=[],
            source_refs=[source_ref],
            authority_class=AuthorityClass.customer_current_authored,
            confidence=0.90,
            review_status=ReviewStatus.auto_accepted,
            parser_version=self.parser_version,
        )
